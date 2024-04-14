from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt


def change_font_and_size(presentation, content_font_size=Pt(32), title_font_size=Pt(36)):
    """
    The function `change_font_and_size` changes the font and size of text in a PowerPoint presentation,
    with different font sizes for content and title slides.

    :param presentation: The presentation parameter is the PowerPoint presentation object that you want
    to modify. It should be an instance of the Presentation class from the python-pptx library
    :param content_font_size: The content_font_size parameter is the font size (in points) that will be
    applied to the content text in each slide of the presentation. The default value is 32 points
    :param title_font_size: The title_font_size parameter is used to specify the font size for the title
    of each slide in the presentation. The default value is set to Pt(36), which represents a font size
    of 36 points. However, you can change this value to any desired font size
    """
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if shape == slide.shapes.title:
                            run.font.name = 'Times New Roman'
                            run.font.size = title_font_size
                            paragraph.alignment = PP_ALIGN.CENTER
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            run.font.name = 'Times New Roman'
                            run.font.size = content_font_size
                            paragraph.alignment = PP_ALIGN.LEFT
                            run.font.color.rgb = RGBColor(0, 0, 0)


def create_presentation(slides_content):
    """
    The function `create_presentation` takes a list of slide content dictionaries, creates a PowerPoint
    presentation, and populates each slide with the corresponding title and content.

    :param slides_content: The `slides_content` parameter is a list of dictionaries, where each
    dictionary represents the content of a slide. Each dictionary should have two keys: 'title' and
    'content'. The value of 'title' should be a string representing the title of the slide, and the
    value of 'content'
    :return: a presentation object.
    """
    presentation = Presentation()
    for slide in slides_content:
        slide_layout = presentation.slide_layouts[1]
        slide_obj = presentation.slides.add_slide(slide_layout)
        title_placeholder = slide_obj.shapes.title
        body_placeholder = slide_obj.placeholders[1]

        title_placeholder.text = slide['title']
        tf = body_placeholder.text_frame
        tf.text = slide['content']

        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(32)

    return presentation


if __name__ == "__main__":
    choice = int(input("Enter \n1 to create a new presentation or\n2 to adjust the size of the presentation: "))
    if choice == 1:
        filename = input("Enter the filename: ") + ".pptx"
        slides_content = [
            {"title": "Title 1", "content": "Content for slide 1"},
            {"title": "Title 2", "content": "Content for slide 2"},
            {"title": "Title 3", "content": "Content for slide 3"}
        ]


        presentation = create_presentation(slides_content)
        change_font_and_size(presentation, Pt(32), Pt(36))
        presentation.save(filename)
        print(f"Presentation saved as {filename}")
    elif choice == 2:
        filename = input("Enter the filename: ") + ".pptx"
        presentation = Presentation(filename)
        change_font_and_size(presentation, Pt(32), Pt(36))
        modified_pptx_file = filename.replace('.pptx', '_modified.pptx')
        presentation.save(modified_pptx_file)
        print(
            f"Font and size changed to Times New Roman and saved as {modified_pptx_file}")
